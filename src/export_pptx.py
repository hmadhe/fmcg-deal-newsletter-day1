"""
Renders the structured newsletter dict into a formatted PowerPoint deck
using python-pptx. One slide per top deal, plus title/summary/other-deals slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BRAND_COLOR = RGBColor(0x1F, 0x4E, 0x79)
BADGE_COLOR = RGBColor(0xC9, 0xDA, 0xEE)
BLANK_LAYOUT_INDEX = 6


def export_newsletter_to_pptx(newsletter, output_path="../output/newsletter.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, newsletter)
    _add_summary_slide(prs, newsletter)

    for deal in newsletter["top_deals"]:
        _add_deal_slide(prs, deal)

    if newsletter["other_deals"]:
        _add_other_deals_slide(prs, newsletter["other_deals"])

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


def _add_slide_heading(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BRAND_COLOR


def _add_title_slide(prs, newsletter):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_COLOR
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)  # send behind any later shapes

    badge = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8), Inches(0.5))
    badge.text_frame.text = newsletter["period"].upper()
    badge.text_frame.paragraphs[0].font.size = Pt(14)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = BADGE_COLOR

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = newsletter["title"]
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(1))
    sub.text_frame.text = (
        f"Generated {newsletter['generated_on']}  |  {newsletter['total_deals_found']} deals tracked this period"
    )
    sub.text_frame.paragraphs[0].font.size = Pt(16)
    sub.text_frame.paragraphs[0].font.color.rgb = BADGE_COLOR


def _add_summary_slide(prs, newsletter):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    _add_slide_heading(slide, "Executive Summary")

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = newsletter["executive_summary"]
    tf.paragraphs[0].font.size = Pt(20)


def _add_deal_slide(prs, deal):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    _add_slide_heading(slide, deal["headline"])

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.7))
    tf = body.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = f"{deal['acquirer']} \u2192 {deal['target']}"
    p0.font.bold = True
    p0.font.size = Pt(22)

    p1 = tf.add_paragraph()
    p1.text = f"{deal['deal_type'].replace('_', ' ').title()}   |   Deal value: {deal['deal_value']}"
    p1.font.size = Pt(16)
    p1.font.color.rgb = BRAND_COLOR

    p2 = tf.add_paragraph()
    p2.text = deal["summary"]
    p2.font.size = Pt(18)

    p3 = tf.add_paragraph()
    n = deal["corroboration_count"]
    p3.text = (
        f"Deal ID: {deal['deal_id']}  |  Sources: {', '.join(deal['sources'])}  |  "
        f"Corroboration: {n} independent source{'s' if n != 1 else ''}  |  "
        f"Confidence: {deal['confidence']}"
    )
    p3.font.size = Pt(12)
    p3.font.italic = True


def _add_other_deals_slide(prs, other_deals):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    _add_slide_heading(slide, "Other Notable Activity")

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, deal in enumerate(other_deals):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = (
            f"\u2022 {deal['acquirer']} \u2192 {deal['target']}: {deal['summary']} "
            f"({', '.join(deal['sources'])})"
        )
        p.font.size = Pt(16)


if __name__ == "__main__":
    import json
    with open("../data/newsletter.json", "r", encoding="utf-8") as f:
        newsletter = json.load(f)
    path = export_newsletter_to_pptx(newsletter)
    print(f"Saved -> {path}")
